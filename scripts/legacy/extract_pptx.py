#!/usr/bin/env python3
from pptx import Presentation
import os

# Load PowerPoint
pptx_path = os.path.expanduser('~/ProjectScraper/Cohesity/PricingCommittee_M365_Feb26.pptx')
prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}\n")

# Extract slides 10-13 (0-indexed: slides 9-12)
for i in [9, 10, 11, 12]:
    if i < len(prs.slides):
        slide = prs.slides[i]
        print(f"{'='*80}")
        print(f"SLIDE {i+1}")
        print(f"{'='*80}")
        
        # Extract all text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text)
                print("-" * 40)
            
            # Check if shape has a table
            if shape.has_table:
                table = shape.table
                print("\n[TABLE FOUND]")
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    print(" | ".join(row_text))
                print("-" * 40)
        
        print("\n")
