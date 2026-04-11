#!/usr/bin/env python3
"""
pptx_writer.py — Durable PPTX write + delivery for Jarvis

Enforces the mandatory durable write protocol:
  1. Atomic write: save to ~/Documents/filename.tmp, verify size, rename to final
  2. Dropbox backup: upload to /Jarvis Reports/ (non-blocking)
  3. Email delivery: delegates to send_file_email.py (3-step validated send)

Usage (as library):
  from pptx_writer import save_pptx_durable, backup_to_dropbox, deliver_pptx

Usage (as CLI — generate a test deck):
  python3 pptx_writer.py --test --subject "Test Deck" --to ericfbrown1@gmail.com

Standing rule (AGENTS.md 2026-04-10):
  All generated PPTX files must use save_pptx_durable() — never prs.save('/tmp/...')
  Always write to ~/Documents/, verify size > 5KB, backup to Dropbox, then send.
"""

import argparse
import hashlib
import os
import shutil
import subprocess
import sys

# ── Constants ────────────────────────────────────────────────────────────────

DOCUMENTS_DIR = os.path.expanduser("~/Documents")
DROPBOX_REPORTS_PATH = "/Jarvis Reports"
SEND_SCRIPT = os.path.expanduser("~/.openclaw/workspace/scripts/send_file_email.py")
DROPBOX_CLI = os.path.expanduser("~/.openclaw/workspace/dropbox-cli.py")
MIN_PPTX_SIZE = 5_000   # bytes — any valid PPTX is larger than this
DEFAULT_TO = "ericfbrown1@gmail.com"
DEFAULT_CC = "Eric.brown@cohesity.com"


# ── Step 1: Durable Atomic Write ─────────────────────────────────────────────

def save_pptx_durable(prs, filename: str) -> tuple[str, str]:
    """
    Save a python-pptx Presentation durably using atomic write + verification.

    Protocol:
      1. Resolve final path to ~/Documents/<filename> (never /tmp/ or relative)
      2. Write to <final_path>.tmp on the same filesystem
      3. Verify temp file exists and size > MIN_PPTX_SIZE
      4. shutil.move() → atomic rename on APFS (same filesystem)
      5. Compute and return SHA256 for integrity tracking

    Args:
      prs:      python-pptx Presentation object
      filename: bare filename, e.g. "my-deck.pptx" (no path — always ~/Documents/)

    Returns:
      (final_path, sha256_hex)

    Raises:
      ValueError  — if filename contains path separators or is unsafe
      RuntimeError — if save or verification fails
    """
    # Safety: reject path separators — filename only, path is always ~/Documents/
    if os.sep in filename or "/" in filename:
        raise ValueError(
            f"filename must be a bare name, not a path. Got: {filename!r}\n"
            f"Correct usage: save_pptx_durable(prs, 'my-deck.pptx')"
        )

    if not filename.lower().endswith(".pptx"):
        filename = filename + ".pptx"

    os.makedirs(DOCUMENTS_DIR, exist_ok=True)
    final_path = os.path.join(DOCUMENTS_DIR, filename)
    tmp_path = final_path + ".tmp"

    print(f"\n── Step 1: Durable Write ───────────────────────────────────")
    print(f"   Target: {final_path}")

    # Clean up any leftover .tmp from a previous failed run
    if os.path.exists(tmp_path):
        os.remove(tmp_path)
        print(f"   Cleaned stale temp file: {tmp_path}")

    try:
        # Write to temp path (same filesystem = atomic rename possible)
        prs.save(tmp_path)

        # Verify temp file
        if not os.path.exists(tmp_path):
            raise RuntimeError(f"prs.save() returned but temp file not found: {tmp_path}")

        size = os.path.getsize(tmp_path)
        if size < MIN_PPTX_SIZE:
            os.remove(tmp_path)
            raise RuntimeError(
                f"PPTX temp file is suspiciously small ({size} bytes < {MIN_PPTX_SIZE} min). "
                f"Likely corrupt or empty presentation. Aborting."
            )

        # Atomic rename: on APFS same-filesystem move is atomic
        shutil.move(tmp_path, final_path)

        # Compute SHA256 for integrity tracking
        with open(final_path, "rb") as fh:
            sha256 = hashlib.sha256(fh.read()).hexdigest()

        print(f"   ✅ Saved: {size:,} bytes, sha256={sha256[:16]}...")
        return final_path, sha256

    except Exception:
        # Always clean up temp on failure
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        raise


# ── Step 2: Dropbox Backup ───────────────────────────────────────────────────

def backup_to_dropbox(local_path: str, dropbox_filename: str = None) -> bool:
    """
    Upload file to Dropbox /Jarvis Reports/ as non-blocking backup.

    Non-blocking: returns False on failure but does NOT raise — email send
    proceeds regardless of Dropbox backup outcome.

    Args:
      local_path:       absolute path to local file
      dropbox_filename: override filename in Dropbox (defaults to basename)

    Returns:
      True on success, False on failure
    """
    if dropbox_filename is None:
        dropbox_filename = os.path.basename(local_path)

    dropbox_path = f"{DROPBOX_REPORTS_PATH}/{dropbox_filename}"

    print(f"\n── Step 2: Dropbox Backup ──────────────────────────────────")
    print(f"   Uploading to: {dropbox_path}")

    result = subprocess.run(
        ["python3", DROPBOX_CLI, "upload", local_path, dropbox_path],
        capture_output=True,
        text=True,
        timeout=60,
    )

    if result.returncode == 0:
        print(f"   ✅ Backed up: {dropbox_path}")
        return True
    else:
        print(f"   ⚠️  Dropbox backup failed (non-blocking — email send continues):")
        print(f"   {result.stderr.strip() or result.stdout.strip()}")
        return False


# ── Step 3: Email Delivery ───────────────────────────────────────────────────

def deliver_pptx(
    prs,
    filename: str,
    subject: str,
    body: str,
    to: str = DEFAULT_TO,
    cc: str = DEFAULT_CC,
) -> dict:
    """
    Full durable delivery pipeline: atomic write → Dropbox backup → validated email send.

    Args:
      prs:      python-pptx Presentation object
      filename: bare filename, e.g. "my-deck.pptx"
      subject:  email subject line
      body:     email body text (footer appended automatically by send script)
      to:       recipient email (default: ericfbrown1@gmail.com)
      cc:       CC email (default: Eric.brown@cohesity.com)

    Returns:
      dict with keys: final_path, sha256, dropbox_ok, email_confirmed
    """
    print("=" * 60)
    print("📎 Jarvis PPTX Durable Delivery Pipeline")
    print("=" * 60)

    # Step 1: Atomic write + verify
    final_path, sha256 = save_pptx_durable(prs, filename)

    # Step 2: Dropbox backup (non-blocking)
    dropbox_ok = backup_to_dropbox(final_path)

    # Step 3: Email send via validated 3-step send script
    print(f"\n── Step 3: Email Delivery ──────────────────────────────────")
    result = subprocess.run(
        [
            "python3", SEND_SCRIPT,
            "--file", final_path,
            "--to", to,
            "--cc", cc,
            "--subject", subject,
            "--body", body,
        ],
        capture_output=False,  # stream output live
        text=True,
    )

    email_ok = result.returncode == 0

    print("\n" + "=" * 60)
    print("📊 Delivery Summary")
    print(f"   File:      {final_path}")
    print(f"   SHA256:    {sha256[:32]}...")
    print(f"   Dropbox:   {'✅ backed up' if dropbox_ok else '⚠️  failed (check manually)'}")
    print(f"   Email:     {'✅ delivered + confirmed' if email_ok else '❌ failed'}")
    print("=" * 60)

    if not email_ok:
        raise RuntimeError(
            f"Email delivery failed. File is safe at: {final_path}\n"
            f"Retry with: python3 {SEND_SCRIPT} --file '{final_path}' "
            f"--to {to} --subject '{subject}' --body '{body}'"
        )

    return {
        "final_path": final_path,
        "sha256": sha256,
        "dropbox_ok": dropbox_ok,
        "email_confirmed": email_ok,
    }


# ── CLI Test Mode ────────────────────────────────────────────────────────────

def _build_test_deck():
    """Build a minimal test PPTX to validate the pipeline."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("❌ python-pptx not installed. Run: pip3 install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Durable Write Protocol — Test Deck"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "Generated by pptx_writer.py — atomic write + Dropbox backup + validated send"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x00, 0xB3, 0x6B)
    p2.space_before = Pt(20)

    return prs


def main():
    parser = argparse.ArgumentParser(
        description="Durable PPTX write + delivery pipeline for Jarvis."
    )
    parser.add_argument("--test", action="store_true", help="Generate a test deck and run full pipeline")
    parser.add_argument("--to", default=DEFAULT_TO, help="Recipient email")
    parser.add_argument("--cc", default=DEFAULT_CC, help="CC email")
    parser.add_argument("--subject", default="Jarvis PPTX Durable Write Test", help="Email subject")
    parser.add_argument("--body", default="Test of the durable PPTX write protocol.", help="Email body")
    parser.add_argument("--filename", default="pptx-durable-write-test.pptx", help="Output filename")
    args = parser.parse_args()

    if args.test:
        print("🔧 Building test deck...")
        prs = _build_test_deck()
        deliver_pptx(prs, args.filename, args.subject, args.body, args.to, args.cc)
    else:
        print("Usage: python3 pptx_writer.py --test")
        print("Or import as library: from pptx_writer import save_pptx_durable, deliver_pptx")
        sys.exit(0)


if __name__ == "__main__":
    main()
